"""
agents/exporter.py — leadership-grade exports: PPTX deck, DOCX report, PDF.

Pure Python (python-pptx / python-docx / matplotlib) so it runs inside the
Streamlit app on any machine. Charts are rendered fresh by matplotlib on a
WHITE background (document convention) with the product palette for series.

Design rules applied: white backgrounds, no accent bars/stripes, strong size
contrast, every chart titled, currency-agnostic figures.
"""
from __future__ import annotations

import io
from datetime import date

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages

AMBER, TEAL, RED, SLATE, INKTXT = "#C97F16", "#1E9E8A", "#D14836", "#5B6B7C", "#1A2430"
TIER_COLORS = {"high": RED, "medium": AMBER, "review": SLATE}


# ---------------------------------------------------------------- chart layer
def _fig(w=8.6, h=4.0):
    fig, ax = plt.subplots(figsize=(w, h), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    ax.grid(axis="both", color="#E3E7EB", linewidth=0.7)
    ax.set_axisbelow(True)
    ax.tick_params(colors=INKTXT, labelsize=9)
    return fig, ax


def _png(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()


def _daily_series(state) -> pd.Series | None:
    cm = state.column_map
    df = state.feature_df if state.feature_df is not None else state.raw_df
    try:
        d = df[[cm["date_column"], cm["target_column"]]].copy()
        d[cm["date_column"]] = pd.to_datetime(d[cm["date_column"]])
        return d.groupby(cm["date_column"])[cm["target_column"]].sum().sort_index()
    except Exception:  # noqa: BLE001
        return None


def render_charts(state) -> dict[str, bytes]:
    """Every chart the documents need, as PNG bytes."""
    charts: dict[str, bytes] = {}
    r = state.results
    series = _daily_series(state)

    fc = r.get("forecasting", {})
    if series is not None and fc.get("forecast"):
        fig, ax = _fig()
        hist = series.tail(120)
        ax.plot(hist.index, hist.values, color=SLATE, lw=1.6, label="Actual")
        fdates = pd.date_range(series.index.max(), periods=len(fc["forecast"]) + 1,
                               freq=pd.infer_freq(series.index) or "D")[1:]
        certified = fc.get("verdict") != "escalate"
        ax.plot([hist.index[-1], *fdates], [hist.values[-1], *fc["forecast"]],
                color=TEAL if certified else AMBER, lw=2.2,
                ls="--" if certified else ":",
                label=f"Forecast · {fc.get('winner_label', '')}"
                      + ("" if certified else " (NOT certified)"))
        ax.axvspan(hist.index[-1], fdates[-1], color=TEAL, alpha=0.05)
        ax.legend(frameon=False, fontsize=9)
        ax.set_title("Revenue — actual and forward view", fontsize=12,
                     color=INKTXT, loc="left", pad=10)
        charts["forecast"] = _png(fig)

    ok = [m for m in fc.get("metrics", [])
          if not m.get("error") and m.get("mape") is not None]
    if ok:
        fig, ax = _fig(8.6, 0.55 * len(ok) + 1.2)
        mdf = pd.DataFrame(ok).sort_values("mape", ascending=True)
        members = (fc.get("ensemble", {}) or {}).get("members") or \
            [fc.get("winner"), fc.get("best_available")]
        colors = [TEAL if m in members else "#C9D2DA" for m in mdf["model"]]
        ax.barh(mdf["model"] + "  ·  " + mdf["family"], mdf["mape"], color=colors)
        for y, v in enumerate(mdf["mape"]):
            ax.text(v, y, f" {v:.2f}%", va="center", fontsize=9, color=INKTXT)
        ax.invert_yaxis()
        ax.set_title("Model bake-off — backtest MAPE (lower is better)",
                     fontsize=12, color=INKTXT, loc="left", pad=10)
        charts["bakeoff"] = _png(fig)

    an = r.get("anomaly", {})
    if series is not None and an.get("flagged"):
        fig, ax = _fig()
        ax.plot(series.index, series.values, color="#C9D2DA", lw=1.1, zorder=1)
        adf = pd.DataFrame([f for f in an["flagged"]
                            if f.get("tier") in ("high", "medium")])
        if len(adf):
            adf["date"] = pd.to_datetime(adf["date"], errors="coerce")
            tcol = state.column_map["target_column"]
            for tier in ("medium", "high"):
                sub = adf[adf["tier"] == tier]
                if len(sub):
                    ax.scatter(sub["date"], sub[tcol], s=14 + sub["votes"] * 7,
                               c=TIER_COLORS[tier], alpha=0.85, zorder=2,
                               label=f"{tier} ({len(sub)})", edgecolors="none")
            ax.legend(frameon=False, fontsize=9)
        ax.set_title("Anomaly constellation — dot size = detector votes",
                     fontsize=12, color=INKTXT, loc="left", pad=10)
        charts["anomaly"] = _png(fig)

    lk = r.get("leakage", {})
    fired = {k: v for k, v in lk.get("rules_fired", {}).items() if v.get("hits")}
    if fired:
        fig, ax = _fig(8.6, 0.6 * len(fired) + 1.2)
        fdf = pd.DataFrame([{"rule": k, "impact": v["impact_total"], "hits": v["hits"]}
                            for k, v in fired.items()]).sort_values("impact")
        ax.barh(fdf["rule"], fdf["impact"], color=AMBER)
        for y, (v, h) in enumerate(zip(fdf["impact"], fdf["hits"])):
            ax.text(v, y, f" {v:,.0f} · {h} hits", va="center", fontsize=9, color=INKTXT)
        ax.set_title("Recoverable impact by leakage rule", fontsize=12,
                     color=INKTXT, loc="left", pad=10)
        charts["leakage"] = _png(fig)

    wf = r.get("whatif", {})
    if wf.get("scenarios"):
        fig, ax = _fig(8.6, 0.6 * len(wf["scenarios"]) + 1.2)
        sdf = pd.DataFrame(wf["scenarios"]).sort_values("delta_vs_baseline")
        ax.barh(sdf["scenario"], sdf["delta_vs_baseline"],
                color=[RED if v < 0 else TEAL for v in sdf["delta_vs_baseline"]])
        for y, (v, pc) in enumerate(zip(sdf["delta_vs_baseline"], sdf["delta_pct"])):
            ax.text(v, y, f" {v:+,.0f} ({pc:+.1f}%)", va="center", fontsize=9,
                    color=INKTXT)
        ax.set_title("What-if scenarios vs baseline", fontsize=12, color=INKTXT,
                     loc="left", pad=10)
        charts["whatif"] = _png(fig)
    return charts


# ---------------------------------------------------------------- shared facts
def _headline(state) -> tuple[str, str]:
    r = state.results
    lk = r.get("leakage", {})
    if lk and "error" not in lk and lk.get("total_impact_estimate"):
        return (f"{lk['total_impact_estimate']:,.0f}",
                "estimated recoverable revenue leakage")
    fc = r.get("forecasting", {})
    return (fc.get("winner_label", "—"), "forecast verdict")


def _key_facts(state) -> list[str]:
    r = state.results
    facts = []
    lk = r.get("leakage", {})
    if lk and "error" not in lk:
        facts.append(f"Revenue leakage: {lk.get('total_impact_estimate', 0):,.0f} "
                     f"across {lk.get('candidate_count', 0)} billing records")
    an = r.get("anomaly", {})
    if an and "error" not in an:
        c = an.get("counts", {})
        facts.append(f"Anomalies: {c.get('high', 0)} high / {c.get('medium', 0)} "
                     f"medium ({len(an.get('detectors_run', []))} detectors voting)")
    fc = r.get("forecasting", {})
    if fc and "error" not in fc:
        wc = fc.get("winner_card") or {}
        mape = f" (MAPE {wc['mape']:.2f}%)" if wc.get("mape") is not None else ""
        facts.append(f"Forecast: {fc.get('verdict')} — {fc.get('winner_label')}{mape}")
    sg = r.get("segment", {})
    if sg and "error" not in sg:
        facts.append(f"Customer base: {sg.get('k')} segments "
                     f"(silhouette {sg.get('silhouette', 0):.2f})")
    return facts


def _actions(state) -> list[dict]:
    rec = state.results.get("recommend", {})
    return rec.get("recommendations", []) if "error" not in rec else []


def _ledger_highlights(state, n=8) -> list[str]:
    keep = ("planner", "critic", "plan_approval", "domain_detection", "narration")
    out = []
    for e in state.ledger.entries:
        if e.stage in keep:
            out.append(f"[{e.stage} · {e.agent}] {e.decision}")
        if len(out) >= n:
            break
    return out


# ---------------------------------------------------------------- PPTX
def build_pptx(state, summary: str, path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    INKC, TEALC, AMBERC = (RGBColor(0x1A, 0x24, 0x30), RGBColor(0x1E, 0x9E, 0x8A),
                           RGBColor(0xC9, 0x7F, 0x16))
    GRAYC = RGBColor(0x5B, 0x6B, 0x7C)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    charts = render_charts(state)

    def _txt(slide, x, y, w, h, text, size, bold=False, color=INKC,
             align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
        return box

    # 1 — title: the headline number IS the slide
    s1 = prs.slides.add_slide(blank)
    num, what = _headline(state)
    _txt(s1, 0.8, 1.6, 11.7, 0.6, "Revenue Reasoning Agent — Findings", 20,
         color=GRAYC)
    _txt(s1, 0.8, 2.3, 11.7, 1.8, num, 88, bold=True, color=AMBERC)
    _txt(s1, 0.8, 4.2, 11.7, 0.7, what, 26, color=INKC)
    _txt(s1, 0.8, 6.6, 11.7, 0.5,
         f"{state.domain.get('name', '')} · run {state.run_id} · {date.today():%d %b %Y}",
         12, color=GRAYC)

    # 2 — key findings
    s2 = prs.slides.add_slide(blank)
    _txt(s2, 0.8, 0.5, 11.7, 0.8, "What the agent found", 36, bold=True)
    box = s2.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, f in enumerate(_key_facts(state)):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = f
        run.font.size, run.font.color.rgb = Pt(20), INKC
        para.space_after = Pt(16)

    # chart slides
    def _chart_slide(title, key, note=None):
        if key not in charts:
            return
        sl = prs.slides.add_slide(blank)
        _txt(sl, 0.8, 0.4, 11.7, 0.8, title, 32, bold=True)
        sl.shapes.add_picture(io.BytesIO(charts[key]), Inches(0.9), Inches(1.4),
                              width=Inches(11.5))
        if note:
            _txt(sl, 0.9, 6.8, 11.5, 0.5, note, 13, color=GRAYC)

    fc = state.results.get("forecasting", {})
    wc = fc.get("winner_card") or {}
    note = None
    if wc:
        note = (f"Winner: {fc.get('winner_label')} — {wc.get('why_selected', '')}"
                )[:220]
    from agents import stage_stories as _SS
    _chart_slide("Forward view", "forecast", note or _SS.verdict_story(state)
                 .replace("**", "")[:220])
    _chart_slide("The bake-off — every model held to account", "bakeoff",
                 (fc.get("ensemble_eval") or {}).get("reason", "")[:220] or None)
    _chart_slide("Where the anomalies sit", "anomaly",
                 _SS.anomaly_story(state).replace("**", "")[:220])
    _chart_slide("Where the money is leaking", "leakage",
                 _SS.leakage_story(state).replace("**", "")[:220])
    _chart_slide("What-if levers", "whatif")

    # actions
    acts = _actions(state)
    if acts:
        s = prs.slides.add_slide(blank)
        _txt(s, 0.8, 0.5, 11.7, 0.8, "Do next", 36, bold=True)
        box = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        for i, a in enumerate(acts[:4]):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = para.add_run()
            run.text = f"{i + 1}.  {a.get('action', '')}"
            run.font.size, run.font.bold, run.font.color.rgb = Pt(20), True, INKC
            sub = tf.add_paragraph()
            sr = sub.add_run()
            sr.text = (f"     effort {a.get('effort', '?')} · confidence "
                       f"{a.get('confidence', 0):.0%} · evidence {a.get('traces_to', '?')}")
            sr.font.size, sr.font.color.rgb = Pt(13), GRAYC
            sub.space_after = Pt(14)

    # how it reasoned
    s = prs.slides.add_slide(blank)
    _txt(s, 0.8, 0.5, 11.7, 0.8, "How the agent reasoned", 36, bold=True)
    box = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(_ledger_highlights(state)):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line[:180]
        run.font.size, run.font.color.rgb = Pt(14), INKC
        para.space_after = Pt(10)
    _txt(s, 0.9, 6.9, 11.5, 0.4,
         "Full decision ledger (JSONL) accompanies this deck.", 12, color=GRAYC)

    prs.save(path)
    return path


# ---------------------------------------------------------------- DOCX
def build_docx(state, summary: str, path: str) -> str:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    charts = render_charts(state)
    doc = Document()
    doc.add_heading("Revenue Reasoning Agent — Run Report", level=0)
    p = doc.add_paragraph(f"Domain: {state.domain.get('name', '?')} · "
                          f"Run {state.run_id} · {date.today():%d %B %Y}")
    p.runs[0].font.color.rgb = RGBColor(0x5B, 0x6B, 0x7C)

    num, what = _headline(state)
    hp = doc.add_paragraph()
    hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r1 = hp.add_run(num + "  ")
    r1.font.size, r1.font.bold = Pt(40), True
    r1.font.color.rgb = RGBColor(0xC9, 0x7F, 0x16)
    r2 = hp.add_run(what)
    r2.font.size = Pt(14)

    doc.add_heading("Executive summary", level=1)
    for para in (summary or "").split("\n\n"):
        text = para.replace("###", "").replace("**", "").replace("*", "").strip()
        if text:
            doc.add_paragraph(text)

    doc.add_heading("Key findings", level=1)
    for f in _key_facts(state):
        doc.add_paragraph(f, style="List Bullet")

    order = [("Forward view", "forecast"), ("Model bake-off", "bakeoff"),
             ("Anomalies", "anomaly"), ("Revenue leakage", "leakage"),
             ("What-if scenarios", "whatif")]
    for title, key in order:
        if key in charts:
            doc.add_heading(title, level=1)
            doc.add_picture(io.BytesIO(charts[key]), width=Inches(6.4))

    fc = state.results.get("forecasting", {})
    if fc.get("roster"):
        doc.add_heading("Model accountability board", level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        for i, h in enumerate(("Model", "Family", "Status", "Why")):
            hdr[i].text = h
        for name, v in fc["roster"].items():
            row = table.add_row().cells
            row[0].text = name
            row[1].text = str(v.get("family", ""))
            row[2].text = str(v.get("status", ""))
            row[3].text = str(v.get("reason", ""))[:160]

    acts = _actions(state)
    if acts:
        doc.add_heading("Recommended actions", level=1)
        for a in acts:
            doc.add_paragraph(
                f"{a.get('action', '')} — effort {a.get('effort', '?')}, "
                f"confidence {a.get('confidence', 0):.0%}, evidence "
                f"{a.get('traces_to', '?')}", style="List Number")

    doc.add_heading("Appendix — decision ledger", level=1)
    for e in state.ledger.entries:
        doc.add_paragraph(f"[{e.stage} · {e.agent}] {e.decision}",
                          style="List Bullet")
    doc.save(path)
    return path


# ---------------------------------------------------------------- PDF
def build_pdf(state, summary: str, path: str) -> str:
    charts = render_charts(state)
    num, what = _headline(state)
    with PdfPages(path) as pdf:
        fig = plt.figure(figsize=(11, 8.5), dpi=150)
        fig.patch.set_facecolor("white")
        fig.text(0.08, 0.80, "Revenue Reasoning Agent — Findings",
                 fontsize=16, color="#5B6B7C")
        fig.text(0.08, 0.58, num, fontsize=54, fontweight="bold", color=AMBER)
        fig.text(0.08, 0.50, what, fontsize=18, color=INKTXT)
        facts = _key_facts(state)
        for i, f in enumerate(facts[:4]):
            fig.text(0.08, 0.38 - i * 0.06, "•  " + f, fontsize=12, color=INKTXT)
        fig.text(0.08, 0.08, f"{state.domain.get('name', '')} · run "
                 f"{state.run_id} · {date.today():%d %b %Y}",
                 fontsize=10, color="#5B6B7C")
        pdf.savefig(fig)
        plt.close(fig)

        import matplotlib.image as mpimg
        for key in ("forecast", "bakeoff", "anomaly", "leakage", "whatif"):
            if key not in charts:
                continue
            img = mpimg.imread(io.BytesIO(charts[key]), format="png")
            fig = plt.figure(figsize=(11, 8.5), dpi=150)
            fig.patch.set_facecolor("white")
            ax = fig.add_axes([0.04, 0.05, 0.92, 0.88])
            ax.imshow(img)
            ax.axis("off")
            pdf.savefig(fig)
            plt.close(fig)
    return path
